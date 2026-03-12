#!/usr/bin/env python3
"""
Generate professional PowerPoint presentations for each chapter of
The Serial Bioinformatician textbook.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Define chapter content
CHAPTERS = {
    1: {
        "title": "Introduction to the Central Dogma",
        "subtitle": "DNA → RNA → Protein",
        "objectives": [
            "Understand the flow of genetic information",
            "Describe the structure and function of DNA",
            "Explain transcription and translation processes",
            "Identify exceptions to the central dogma",
            "Recognize the importance in bioinformatics"
        ],
        "sections": [
            ("The Central Dogma", [
                "Francis Crick proposed in 1958",
                "DNA is the repository of genetic information",
                "Information flows: DNA → RNA → Protein",
                "One gene typically encodes one protein",
                "Foundation of molecular biology"
            ]),
            ("DNA Structure", [
                "Double helix with antiparallel strands",
                "Sugar-phosphate backbone",
                "Four bases: A, T, G, C (A-T, G-C pairing)",
                "5' to 3' directionality",
                "Genetic code stored in sequence"
            ]),
            ("Transcription", [
                "DNA → mRNA in the nucleus",
                "RNA Polymerase II reads template strand",
                "5' cap and 3' poly-A tail added",
                "Splicing removes introns",
                "mRNA exported to cytoplasm"
            ]),
            ("Translation", [
                "mRNA → Protein at ribosome",
                "tRNA brings amino acids",
                "Codon-anticodon base pairing",
                "Start (AUG) and stop codons (UAA, UAG, UGA)",
                "Polypeptide chain folds into functional protein"
            ])
        ],
        "summary": [
            "Central dogma: DNA → RNA → Protein",
            "Transcription occurs in nucleus",
            "Translation occurs in cytoplasm",
            "Exceptions: reverse transcriptase, RNA viruses",
            "Fundamental to understanding genomics"
        ]
    },
    2: {
        "title": "The Genome and its Variations",
        "subtitle": "Understanding Genetic Diversity",
        "objectives": [
            "Define genome structure and organization",
            "Classify types of genetic variants",
            "Understand mutation mechanisms",
            "Interpret variant notation (HGVS)",
            "Recognize clinical significance of variants"
        ],
        "sections": [
            ("Genome Organization", [
                "Human genome: ~3 billion base pairs",
                "~20,000 protein-coding genes",
                "Exons, introns, regulatory regions",
                "Repetitive elements (~45%)",
                "Reference genomes: GRCh38, T2T-CHM13"
            ]),
            ("Types of Variants", [
                "SNVs: Single nucleotide variants",
                "Indels: Insertions and deletions",
                "CNVs: Copy number variations",
                "Structural variants: inversions, translocations",
                "Repeat expansions: trinucleotide repeats"
            ]),
            ("Variant Annotation", [
                "HGVS nomenclature standard",
                "Genomic vs cDNA vs protein notation",
                "c.123A>G, p.Val600Glu",
                "Consequence prediction: missense, nonsense",
                "Tools: VEP, ANNOVAR, SnpEff"
            ]),
            ("Clinical Significance", [
                "ACMG classification guidelines",
                "Pathogenic, Likely Pathogenic, VUS, Benign",
                "gnomAD population frequencies",
                "ClinVar database",
                "Evidence-based variant interpretation"
            ])
        ],
        "summary": [
            "Genomes contain coding and non-coding regions",
            "Variants range from SNVs to chromosomal rearrangements",
            "HGVS provides standardized nomenclature",
            "ACMG guidelines for clinical interpretation",
            "Population databases essential for filtering"
        ]
    },
    3: {
        "title": "Proteins: The Functional Units",
        "subtitle": "From Sequence to Structure to Function",
        "objectives": [
            "Understand protein structure hierarchy",
            "Describe protein folding principles",
            "Explain structure-function relationships",
            "Use AlphaFold for structure prediction",
            "Interpret protein structural data"
        ],
        "sections": [
            ("Protein Structure Hierarchy", [
                "Primary: amino acid sequence",
                "Secondary: α-helices, β-sheets",
                "Tertiary: 3D fold",
                "Quaternary: multi-subunit complexes",
                "Structure determines function"
            ]),
            ("Protein Folding", [
                "Thermodynamically favored state",
                "Hydrophobic core, hydrophilic surface",
                "Chaperones assist folding",
                "Misfolding: Alzheimer's, Parkinson's",
                "Levinthal's paradox"
            ]),
            ("AlphaFold Revolution", [
                "DeepMind's breakthrough (2020)",
                "Predicts 3D structure from sequence",
                ">200 million structures predicted",
                "pLDDT confidence scores",
                "ColabFold for fast inference"
            ]),
            ("Structural Analysis Tools", [
                "PyMOL, ChimeraX visualization",
                "RCSB PDB database",
                "UniProt sequence features",
                "Pfam/InterPro domains",
                "Docking and dynamics simulations"
            ])
        ],
        "summary": [
            "Proteins have four levels of structure",
            "Sequence determines structure determines function",
            "AlphaFold revolutionized structure prediction",
            "Structural databases: PDB, AlphaFold DB",
            "Visualization tools enable analysis"
        ]
    },
    4: {
        "title": "Command Line for Biologists",
        "subtitle": "Essential Unix/Linux Skills",
        "objectives": [
            "Navigate the filesystem via command line",
            "Manipulate text files with core utilities",
            "Write basic shell scripts",
            "Use pipes and redirects effectively",
            "Process biological data with CLI tools"
        ],
        "sections": [
            ("Basic Navigation", [
                "pwd, ls, cd - directory operations",
                "mkdir, rmdir, rm - file management",
                "cp, mv - copy and move",
                "Absolute vs relative paths",
                "Hidden files and permissions"
            ]),
            ("Text Processing", [
                "cat, head, tail, less - viewing files",
                "grep - pattern searching",
                "sed, awk - text transformation",
                "sort, uniq, cut, paste",
                "wc - word/line counting"
            ]),
            ("Pipes and Redirects", [
                "| pipe: command1 | command2",
                "> redirect output to file",
                ">> append to file",
                "< input from file",
                "2>&1 redirect stderr"
            ]),
            ("Bioinformatics CLI Tools", [
                "samtools - BAM file manipulation",
                "bedtools - genomic intervals",
                "bcftools - VCF processing",
                "seqkit - FASTA/FASTQ handling",
                "GNU parallel - parallelization"
            ])
        ],
        "summary": [
            "Command line is essential for bioinformatics",
            "Text processing tools for data manipulation",
            "Pipes chain commands together",
            "Shell scripts automate workflows",
            "Domain-specific tools for genomics data"
        ]
    },
    5: {
        "title": "Python for Bioinformatics",
        "subtitle": "Programming Biological Analysis",
        "objectives": [
            "Write Python scripts for data analysis",
            "Use Biopython for sequence manipulation",
            "Work with pandas DataFrames",
            "Create visualizations with matplotlib/seaborn",
            "Build reproducible analysis pipelines"
        ],
        "sections": [
            ("Python Fundamentals", [
                "Variables, data types, operators",
                "Lists, dictionaries, sets",
                "Control flow: if, for, while",
                "Functions and modules",
                "Virtual environments"
            ]),
            ("Biopython Library", [
                "SeqIO: reading/writing sequences",
                "Seq objects and operations",
                "BLAST searches programmatically",
                "PDB structure parsing",
                "Entrez: NCBI API access"
            ]),
            ("Data Analysis Stack", [
                "pandas: DataFrames and Series",
                "numpy: numerical computing",
                "scipy: scientific functions",
                "scikit-learn: machine learning",
                "statsmodels: statistical tests"
            ]),
            ("Visualization", [
                "matplotlib: base plotting",
                "seaborn: statistical graphics",
                "plotly: interactive plots",
                "Specialized: logomaker, dna_features_viewer",
                "Publication-quality figures"
            ])
        ],
        "summary": [
            "Python is the premier bioinformatics language",
            "Biopython handles sequences and databases",
            "pandas/numpy for data manipulation",
            "Rich visualization ecosystem",
            "Jupyter notebooks for exploration"
        ]
    },
    6: {
        "title": "Navigating Biological Databases",
        "subtitle": "Finding and Retrieving Data",
        "objectives": [
            "Navigate major biological databases",
            "Retrieve sequences and annotations",
            "Use APIs for programmatic access",
            "Understand data formats and standards",
            "Integrate data from multiple sources"
        ],
        "sections": [
            ("Sequence Databases", [
                "NCBI GenBank: nucleotide sequences",
                "UniProt: protein sequences + annotations",
                "Ensembl: genome browser",
                "RefSeq: curated reference sequences",
                "EMBL-EBI resources"
            ]),
            ("Structure Databases", [
                "RCSB PDB: experimental structures",
                "AlphaFold DB: predicted structures",
                "UniProt structure annotations",
                "PDBe: European resource",
                "ModelArchive: computational models"
            ]),
            ("Functional Databases", [
                "Gene Ontology (GO)",
                "KEGG: pathways and metabolism",
                "Reactome: biological pathways",
                "InterPro: protein families",
                "STRING: protein interactions"
            ]),
            ("API Access", [
                "REST APIs: HTTP requests",
                "Biopython Entrez module",
                "biomaRt (R/Bioconductor)",
                "Rate limiting and best practices",
                "Local database installations"
            ])
        ],
        "summary": [
            "Curated databases are essential resources",
            "NCBI, UniProt, Ensembl are key portals",
            "Programmatic access via APIs",
            "Multiple databases for comprehensive analysis",
            "Data integration is a key skill"
        ]
    },
    7: {
        "title": "Sequence Alignment",
        "subtitle": "Comparing Biological Sequences",
        "objectives": [
            "Understand alignment algorithms",
            "Perform pairwise and multiple alignments",
            "Interpret alignment statistics",
            "Choose appropriate alignment tools",
            "Apply alignment to biological questions"
        ],
        "sections": [
            ("Alignment Fundamentals", [
                "Scoring matrices: BLOSUM, PAM",
                "Gap penalties: opening vs extension",
                "Global vs local alignment",
                "Dynamic programming approach",
                "Heuristic methods for speed"
            ]),
            ("Pairwise Alignment", [
                "Needleman-Wunsch (global)",
                "Smith-Waterman (local)",
                "BLAST: fast local alignment",
                "E-value interpretation",
                "Sequence identity vs similarity"
            ]),
            ("Multiple Sequence Alignment", [
                "MUSCLE, MAFFT, ClustalOmega",
                "Progressive alignment strategy",
                "Iterative refinement",
                "Visualization: Jalview, AliView",
                "Quality assessment"
            ]),
            ("Read Alignment", [
                "BWA-MEM2, Bowtie2 (short reads)",
                "minimap2 (long reads)",
                "STAR (RNA-seq)",
                "SAM/BAM format",
                "Alignment QC metrics"
            ])
        ],
        "summary": [
            "Alignment reveals evolutionary relationships",
            "Choose algorithm based on use case",
            "BLAST for database searches",
            "MSA essential for phylogenetics",
            "Read aligners for NGS data"
        ]
    },
    8: {
        "title": "Phylogenetics",
        "subtitle": "Reconstructing Evolutionary Relationships",
        "objectives": [
            "Build phylogenetic trees from sequences",
            "Understand tree-building methods",
            "Interpret evolutionary relationships",
            "Assess tree confidence (bootstrap)",
            "Apply phylogenetics to research questions"
        ],
        "sections": [
            ("Phylogenetic Basics", [
                "Trees represent evolutionary history",
                "Rooted vs unrooted trees",
                "Branch lengths: evolutionary distance",
                "Clades and common ancestors",
                "Homology vs analogy"
            ]),
            ("Tree-Building Methods", [
                "Distance-based: NJ, UPGMA",
                "Maximum Parsimony",
                "Maximum Likelihood: IQ-TREE, RAxML",
                "Bayesian: MrBayes, BEAST",
                "Model selection: ModelFinder"
            ]),
            ("Tree Confidence", [
                "Bootstrap analysis",
                "SH-aLRT support values",
                "Posterior probabilities (Bayesian)",
                "UFBoot ultrafast bootstrap",
                "Interpreting support values"
            ]),
            ("Specialized Applications", [
                "Molecular clocks and dating",
                "Ancestral state reconstruction",
                "Gene tree vs species tree",
                "Viral phylogenomics (NextStrain)",
                "Outbreak investigation"
            ])
        ],
        "summary": [
            "Phylogenetics reveals evolutionary history",
            "ML and Bayesian methods most accurate",
            "Bootstrap values assess confidence",
            "IQ-TREE is a modern standard tool",
            "Applications: evolution, epidemiology"
        ]
    },
    9: {
        "title": "Genome Assembly and Annotation",
        "subtitle": "Building and Interpreting Genomes",
        "objectives": [
            "Understand de novo assembly strategies",
            "Compare short-read and long-read assembly",
            "Perform assembly quality assessment",
            "Annotate genes and functional elements",
            "Interpret assembly statistics"
        ],
        "sections": [
            ("Assembly Fundamentals", [
                "De novo vs reference-guided",
                "De Bruijn graphs (short reads)",
                "Overlap-layout-consensus (long reads)",
                "Contigs, scaffolds, chromosomes",
                "Hybrid assembly approaches"
            ]),
            ("Assembly Tools", [
                "SPAdes (short reads)",
                "Flye, Canu, Hifiasm (long reads)",
                "Verkko, hifiasm (HiFi + ONT)",
                "Polishing: Pilon, Racon",
                "T2T assembly approaches"
            ]),
            ("Quality Assessment", [
                "N50, L50 statistics",
                "QUAST metrics",
                "BUSCO completeness",
                "k-mer analysis: Merqury",
                "Contamination screening"
            ]),
            ("Genome Annotation", [
                "Gene prediction: Augustus, BRAKER",
                "Functional annotation: InterProScan",
                "RNA-seq guided annotation",
                "Repeat masking: RepeatMasker",
                "NCBI Prokaryotic Annotation Pipeline"
            ])
        ],
        "summary": [
            "Long reads enable complete assemblies",
            "N50 and BUSCO assess quality",
            "Hybrid approaches combine technologies",
            "Annotation combines ab initio and evidence",
            "T2T genomes are the new standard"
        ]
    },
    10: {
        "title": "Introduction to NGS",
        "subtitle": "Next-Generation Sequencing Technologies",
        "objectives": [
            "Understand sequencing technologies",
            "Perform quality control on raw data",
            "Process FASTQ files appropriately",
            "Choose sequencing strategy for experiments",
            "Design NGS experiments effectively"
        ],
        "sections": [
            ("Sequencing Platforms", [
                "Illumina: short-read, high accuracy",
                "PacBio HiFi: long-read, high accuracy",
                "Oxford Nanopore: ultra-long reads",
                "Element Biosciences: emerging platform",
                "MGI/BGI: high throughput"
            ]),
            ("Raw Data Processing", [
                "FASTQ format structure",
                "Quality scores (Phred encoding)",
                "FastQC quality assessment",
                "MultiQC report aggregation",
                "Adapter trimming: fastp, cutadapt"
            ]),
            ("Experimental Design", [
                "Coverage depth requirements",
                "Read length considerations",
                "Paired-end vs single-end",
                "Multiplexing and barcoding",
                "Cost vs quality trade-offs"
            ]),
            ("Workflow Management", [
                "Snakemake workflows",
                "Nextflow pipelines",
                "nf-core community pipelines",
                "Containerization: Docker, Singularity",
                "Reproducibility best practices"
            ])
        ],
        "summary": [
            "Multiple sequencing platforms available",
            "QC is essential before analysis",
            "Match technology to biological question",
            "Workflow managers ensure reproducibility",
            "nf-core provides production-ready pipelines"
        ]
    },
    11: {
        "title": "Transcriptomics",
        "subtitle": "RNA-Seq Analysis from Raw Reads to Insights",
        "objectives": [
            "Process RNA-seq data end-to-end",
            "Perform differential expression analysis",
            "Interpret fold changes and statistics",
            "Conduct pathway enrichment analysis",
            "Visualize transcriptomic results"
        ],
        "sections": [
            ("RNA-Seq Workflow", [
                "QC → Alignment → Quantification",
                "STAR, HISAT2 aligners",
                "Salmon, Kallisto pseudo-alignment",
                "Gene vs transcript quantification",
                "Normalization: TPM, RPKM, counts"
            ]),
            ("Differential Expression", [
                "DESeq2, edgeR, limma-voom",
                "Statistical modeling of counts",
                "Log2 fold change interpretation",
                "Adjusted p-values (FDR)",
                "Volcano plots and MA plots"
            ]),
            ("Pathway Analysis", [
                "Gene Ontology enrichment",
                "KEGG pathway analysis",
                "GSEA and fgsea",
                "Over-representation analysis",
                "ClusterProfiler package"
            ]),
            ("Advanced Topics", [
                "Single-cell RNA-seq (Seurat, Scanpy)",
                "Alternative splicing analysis",
                "Long-read transcriptomics",
                "Spatial transcriptomics",
                "Multi-sample integration"
            ])
        ],
        "summary": [
            "RNA-seq reveals gene expression patterns",
            "DESeq2 is the standard DE method",
            "Always correct for multiple testing",
            "GO/pathway analysis provides context",
            "Single-cell adds cellular resolution"
        ]
    },
    12: {
        "title": "Proteomics and Metabolomics",
        "subtitle": "Beyond the Genome",
        "objectives": [
            "Understand mass spectrometry principles",
            "Analyze proteomics data",
            "Identify and quantify metabolites",
            "Integrate omics data types",
            "Interpret multi-omics results"
        ],
        "sections": [
            ("Mass Spectrometry Basics", [
                "Ionization: ESI, MALDI",
                "Mass analyzers: Orbitrap, TOF",
                "MS/MS fragmentation",
                "DDA vs DIA acquisition",
                "Data formats: mzML, raw"
            ]),
            ("Proteomics Analysis", [
                "Database search: MaxQuant, MSFragger",
                "Label-free quantification",
                "TMT/iTRAQ labeled quantification",
                "Statistical analysis: MSstats, limma",
                "PTM identification"
            ]),
            ("Metabolomics", [
                "Targeted vs untargeted approaches",
                "Feature extraction: XCMS, MZmine",
                "Metabolite identification",
                "Pathway mapping: MetaboAnalyst",
                "Lipidomics specialization"
            ]),
            ("Multi-Omics Integration", [
                "Data normalization across omics",
                "MOFA: multi-omics factor analysis",
                "Network-based integration",
                "Pathway-level integration",
                "Visualization strategies"
            ])
        ],
        "summary": [
            "MS enables protein and metabolite profiling",
            "DIA provides comprehensive coverage",
            "MaxQuant/MSFragger for proteomics search",
            "MetaboAnalyst for metabolomics analysis",
            "Multi-omics reveals systems biology"
        ]
    },
    13: {
        "title": "Microbiomics",
        "subtitle": "Analyzing Microbial Communities",
        "objectives": [
            "Characterize microbial diversity",
            "Perform 16S rRNA gene analysis",
            "Conduct metagenomic analyses",
            "Interpret diversity metrics",
            "Apply microbiome analysis to research"
        ],
        "sections": [
            ("Microbiome Basics", [
                "16S rRNA gene marker",
                "Amplicon sequencing approach",
                "Shotgun metagenomics",
                "Alpha and beta diversity",
                "Community composition"
            ]),
            ("16S Analysis", [
                "QIIME2 workflow",
                "DADA2 for ASV generation",
                "Taxonomic classification",
                "Diversity metrics: Shannon, Chao1",
                "Ordination: PCoA, NMDS"
            ]),
            ("Shotgun Metagenomics", [
                "Assembly: MEGAHIT, metaSPAdes",
                "Binning: MetaBAT2, CONCOCT",
                "MAG quality: CheckM",
                "Functional profiling: HUMAnN",
                "Taxonomic profiling: Kraken2, MetaPhlAn"
            ]),
            ("Statistical Analysis", [
                "Differential abundance: ANCOM, MaAsLin2",
                "Compositional data considerations",
                "Ecological statistics",
                "Machine learning for classification",
                "Longitudinal analysis"
            ])
        ],
        "summary": [
            "Microbiomes are complex communities",
            "DADA2/ASVs are current best practice",
            "MAGs from metagenomes expand discovery",
            "Compositional nature requires special stats",
            "Function often more relevant than taxonomy"
        ]
    },
    14: {
        "title": "Structural Bioinformatics",
        "subtitle": "3D Structures and Molecular Interactions",
        "objectives": [
            "Analyze protein and nucleic acid structures",
            "Perform structure prediction and modeling",
            "Conduct molecular docking",
            "Use visualization tools effectively",
            "Interpret structure-function relationships"
        ],
        "sections": [
            ("Structural Data", [
                "PDB format and mmCIF",
                "RCSB PDB database",
                "AlphaFold Database",
                "Structure quality metrics",
                "Resolution and R-factor"
            ]),
            ("Structure Prediction", [
                "AlphaFold2 and ColabFold",
                "ESMFold for fast prediction",
                "RoseTTAFold",
                "Homology modeling principles",
                "pLDDT confidence interpretation"
            ]),
            ("Molecular Docking", [
                "Protein-ligand docking: AutoDock",
                "Protein-protein docking: HADDOCK",
                "Virtual screening",
                "Binding site prediction",
                "Scoring functions"
            ]),
            ("Visualization and Analysis", [
                "PyMOL for publication figures",
                "ChimeraX for analysis",
                "Structure alignment: TM-align",
                "Surface analysis",
                "Molecular dynamics setup"
            ])
        ],
        "summary": [
            "Structures reveal molecular mechanisms",
            "AlphaFold democratized structure prediction",
            "Docking predicts molecular interactions",
            "Visualization essential for interpretation",
            "Structure guides drug design"
        ]
    },
    15: {
        "title": "Biological Networks",
        "subtitle": "Systems-Level Understanding",
        "objectives": [
            "Construct biological networks",
            "Analyze network topology",
            "Identify functional modules",
            "Integrate multiple data types",
            "Interpret network analysis results"
        ],
        "sections": [
            ("Network Types", [
                "Protein-protein interaction (PPI)",
                "Gene regulatory networks",
                "Metabolic networks",
                "Co-expression networks",
                "Signaling networks"
            ]),
            ("Network Analysis", [
                "Degree distribution",
                "Centrality measures",
                "Clustering coefficient",
                "Network motifs",
                "Scale-free and small-world properties"
            ]),
            ("Module Detection", [
                "Community detection algorithms",
                "WGCNA for co-expression",
                "Louvain/Leiden clustering",
                "Functional module enrichment",
                "Hub gene identification"
            ]),
            ("Tools and Visualization", [
                "Cytoscape visualization",
                "STRING database",
                "igraph (R/Python)",
                "NetworkX (Python)",
                "Network propagation methods"
            ])
        ],
        "summary": [
            "Networks capture biological relationships",
            "Topology reveals important nodes",
            "Modules represent functional units",
            "WGCNA powerful for expression data",
            "Integration enhances biological meaning"
        ]
    },
    16: {
        "title": "GWAS and Population Genetics",
        "subtitle": "Genetic Association Studies",
        "objectives": [
            "Design and analyze GWAS studies",
            "Understand population structure",
            "Perform association testing",
            "Interpret GWAS results",
            "Apply polygenic risk scores"
        ],
        "sections": [
            ("GWAS Fundamentals", [
                "Common disease common variant model",
                "SNP arrays vs sequencing",
                "Case-control and quantitative designs",
                "Genome-wide significance threshold",
                "Effect sizes and odds ratios"
            ]),
            ("Quality Control", [
                "Sample QC: relatedness, sex check",
                "Variant QC: call rate, HWE",
                "Population stratification",
                "Principal components correction",
                "PLINK for QC"
            ]),
            ("Association Testing", [
                "Logistic/linear regression",
                "Mixed models: BOLT-LMM, SAIGE",
                "Rare variant tests: SKAT, burden",
                "Multi-ancestry meta-analysis",
                "Fine-mapping: FINEMAP, SuSiE"
            ]),
            ("Interpretation", [
                "Manhattan and QQ plots",
                "Functional annotation",
                "eQTL colocalization",
                "Polygenic risk scores",
                "Mendelian randomization"
            ])
        ],
        "summary": [
            "GWAS identifies disease associations",
            "QC and population structure critical",
            "5×10⁻⁸ genome-wide significance",
            "Fine-mapping identifies causal variants",
            "PRS enables genetic prediction"
        ]
    },
    17: {
        "title": "Co-expression Network Analysis",
        "subtitle": "WGCNA and Module Detection",
        "objectives": [
            "Build weighted co-expression networks",
            "Identify gene modules",
            "Correlate modules with traits",
            "Find hub genes",
            "Interpret module function"
        ],
        "sections": [
            ("WGCNA Overview", [
                "Weighted gene co-expression networks",
                "Soft thresholding power",
                "Scale-free topology criterion",
                "TOM: Topological Overlap Matrix",
                "Module detection via cutting"
            ]),
            ("Network Construction", [
                "Expression data filtering",
                "Sample clustering for outliers",
                "Power selection",
                "Adjacency and TOM calculation",
                "Module eigengenes"
            ]),
            ("Module-Trait Analysis", [
                "Module-trait correlations",
                "Gene significance and membership",
                "Hub gene identification",
                "Intramodular connectivity",
                "Module preservation analysis"
            ]),
            ("Interpretation", [
                "Gene ontology enrichment",
                "Module visualization",
                "Cross-species preservation",
                "Integration with other data",
                "Biological validation approaches"
            ])
        ],
        "summary": [
            "WGCNA groups co-expressed genes",
            "Modules often represent pathways",
            "Hub genes are key regulators",
            "Module-trait correlations prioritize targets",
            "Powerful for large expression datasets"
        ]
    },
    18: {
        "title": "Interactive Visualization",
        "subtitle": "Shiny Dashboards for Bioinformatics",
        "objectives": [
            "Build interactive R Shiny applications",
            "Design user-friendly interfaces",
            "Deploy dashboards for data exploration",
            "Create publication-ready visualizations",
            "Enable non-programmers to explore data"
        ],
        "sections": [
            ("Shiny Fundamentals", [
                "UI and Server architecture",
                "Reactive programming model",
                "Input and output widgets",
                "Layouts: fluidPage, sidebarLayout",
                "shinydashboard package"
            ]),
            ("Interactive Plots", [
                "plotly for interactivity",
                "ggplot2 + plotly integration",
                "DT for interactive tables",
                "Brushing and linking",
                "Custom JavaScript widgets"
            ]),
            ("Bioinformatics Dashboards", [
                "Gene expression explorers",
                "WGCNA module browsers",
                "Variant filtering interfaces",
                "Pathway visualization",
                "Multi-omics integration views"
            ]),
            ("Deployment", [
                "shinyapps.io hosting",
                "Shiny Server deployment",
                "Docker containerization",
                "Performance optimization",
                "User authentication"
            ])
        ],
        "summary": [
            "Shiny makes R analyses accessible",
            "Reactive model enables dynamic updates",
            "plotly adds rich interactivity",
            "Dashboards democratize data exploration",
            "Multiple deployment options available"
        ]
    },
    19: {
        "title": "Statistics for Bioinformatics",
        "subtitle": "Essential Statistical Methods",
        "objectives": [
            "Apply appropriate statistical tests",
            "Handle multiple testing correction",
            "Perform dimensionality reduction",
            "Interpret statistical results correctly",
            "Avoid common statistical pitfalls"
        ],
        "sections": [
            ("Hypothesis Testing", [
                "t-tests, ANOVA, chi-square",
                "Non-parametric alternatives",
                "Effect sizes and power",
                "Confidence intervals",
                "P-value interpretation"
            ]),
            ("Multiple Testing", [
                "Family-wise error rate",
                "Bonferroni correction",
                "False Discovery Rate (FDR)",
                "Benjamini-Hochberg procedure",
                "q-values"
            ]),
            ("Dimensionality Reduction", [
                "PCA: Principal Component Analysis",
                "t-SNE visualization",
                "UMAP for single-cell",
                "Feature selection methods",
                "Variance explained"
            ]),
            ("Specialized Methods", [
                "Mixed effects models",
                "Survival analysis",
                "Bayesian statistics",
                "Bootstrapping and permutation",
                "Machine learning validation"
            ])
        ],
        "summary": [
            "Match test to data type and design",
            "Always correct for multiple testing",
            "PCA reduces dimensions while preserving variance",
            "Effect sizes matter more than p-values",
            "Proper validation prevents overfitting"
        ]
    },
    20: {
        "title": "Multiomics Integration",
        "subtitle": "Combining Omics for Systems Biology",
        "objectives": [
            "Integrate multiple omics data types",
            "Apply multi-omics analysis methods",
            "Interpret integrated results",
            "Handle technical challenges",
            "Design multi-omics experiments"
        ],
        "sections": [
            ("Integration Strategies", [
                "Early, intermediate, late fusion",
                "Supervised vs unsupervised",
                "Network-based integration",
                "Pathway-level integration",
                "Knowledge-guided approaches"
            ]),
            ("Methods and Tools", [
                "MOFA/MOFA+: factor analysis",
                "mixOmics: multivariate methods",
                "SNF: similarity network fusion",
                "MINT: multi-study integration",
                "iCluster and PARADIGM"
            ]),
            ("Technical Considerations", [
                "Batch effect correction",
                "Missing data handling",
                "Cross-platform normalization",
                "Sample size requirements",
                "Validation strategies"
            ]),
            ("Biological Applications", [
                "Cancer subtyping",
                "Drug response prediction",
                "Biomarker discovery",
                "Mechanism elucidation",
                "Personalized medicine"
            ])
        ],
        "summary": [
            "Multi-omics provides comprehensive view",
            "Integration reveals emergent patterns",
            "MOFA is powerful for unsupervised analysis",
            "Technical challenges require careful handling",
            "Applications across disease research"
        ]
    }
}

# Color scheme
COLORS = {
    "primary": RGBColor(21, 101, 192),      # #1565C0
    "secondary": RGBColor(92, 107, 192),    # #5C6BC0
    "accent": RGBColor(0, 137, 123),         # #00897B
    "dark": RGBColor(38, 50, 56),            # #263238
    "light": RGBColor(236, 239, 241),        # #ECEFF1
    "white": RGBColor(255, 255, 255),
    "highlight": RGBColor(255, 152, 0),      # #FF9800
}

def create_title_slide(prs, chapter_num, chapter_data):
    """Create the title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary"]
    shape.line.fill.background()
    
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5), Inches(10), Inches(0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["highlight"]
    bar.line.fill.background()
    
    # Chapter number
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Chapter {chapter_num}"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS["light"]
    p.font.bold = False
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = chapter_data["title"]
    p.font.size = Pt(44)
    p.font.color.rgb = COLORS["white"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = chapter_data["subtitle"]
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS["light"]
    p.font.bold = False
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "The Serial Bioinformatician"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS["light"]
    p.alignment = PP_ALIGN.CENTER

def create_objectives_slide(prs, chapter_data):
    """Create learning objectives slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["primary"]
    bar.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS["white"]
    p.font.bold = True
    
    # Objectives list
    for i, obj in enumerate(chapter_data["objectives"]):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8 + i * 0.9), Inches(8.5), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"✓  {obj}"
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS["dark"]

def create_content_slide(prs, section_title, bullets):
    """Create a content slide with bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["secondary"]
    bar.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS["white"]
    p.font.bold = True
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["accent"]
    line.line.fill.background()
    
    # Bullets
    for i, bullet in enumerate(bullets):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2 + i * 0.95), Inches(8.5), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"•  {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS["dark"]

def create_summary_slide(prs, chapter_data):
    """Create summary slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["accent"]
    bar.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Takeaways"
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS["white"]
    p.font.bold = True
    
    # Summary points
    for i, point in enumerate(chapter_data["summary"]):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8 + i * 0.95), Inches(8.5), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"→  {point}"
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS["dark"]
        p.font.bold = True

def create_closing_slide(prs, chapter_num):
    """Create closing slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["dark"]
    shape.line.fill.background()
    
    # Thank you
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(48)
    p.font.color.rgb = COLORS["white"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Next chapter hint
    if chapter_num < 20:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Next: Chapter {chapter_num + 1}"
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS["secondary"]
        p.alignment = PP_ALIGN.CENTER
    
    # Footer
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "The Serial Bioinformatician | © 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS["light"]
    p.alignment = PP_ALIGN.CENTER

def create_chapter_presentation(chapter_num, chapter_data, output_dir):
    """Create a complete presentation for a chapter."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    create_title_slide(prs, chapter_num, chapter_data)
    
    # Objectives slide
    create_objectives_slide(prs, chapter_data)
    
    # Content slides
    for section_title, bullets in chapter_data["sections"]:
        create_content_slide(prs, section_title, bullets)
    
    # Summary slide
    create_summary_slide(prs, chapter_data)
    
    # Closing slide
    create_closing_slide(prs, chapter_num)
    
    # Save
    filename = f"chapter-{chapter_num:02d}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    print(f"Created: {filepath}")

def main():
    output_dir = "docs/slides"
    os.makedirs(output_dir, exist_ok=True)
    
    for chapter_num, chapter_data in CHAPTERS.items():
        create_chapter_presentation(chapter_num, chapter_data, output_dir)
    
    print("\n✓ All presentations created successfully!")

if __name__ == "__main__":
    main()
